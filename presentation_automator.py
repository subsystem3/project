import os
import shutil
import subprocess

from gtts import gTTS
from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips
from pdf2image import convert_from_path
from pptx import Presentation


class PresentationAutomator:
    """
    A class to automate the creation of a video presentation from a PowerPoint file.
    """

    def __init__(self, pptx_filename: str, output_file: str):
        self.pptx_filename = pptx_filename
        self.pdf_filename = pptx_filename.replace(".pptx", ".pdf")
        self.output_file = output_file
        self.presentation = Presentation(pptx_filename)
        self.slides = self.presentation.slides
        self.voiceover_texts = [
            slide.notes_slide.notes_text_frame.text for slide in self.slides
        ]

    def convert_to_pdf(self):
        """Converts the .pptx file to a .pdf file using LibreOffice."""
        cmd = f"libreoffice --headless --convert-to pdf {self.pptx_filename}"
        subprocess.run(cmd, shell=True, check=True)

    def create_videos(self):
        """Creates a video for each slide with a voiceover."""
        videos = []
        assets_dir = "assets"
        if os.path.exists(assets_dir):
            shutil.rmtree(assets_dir)
        os.makedirs(assets_dir, exist_ok=True)
        for i, slide in enumerate(self.slides):
            text = self.voiceover_texts[i]
            images = convert_from_path(self.pdf_filename, dpi=200)
            images[i].save(f"{assets_dir}/slide_{i}.png", "PNG")
            voice = gTTS(text=text, lang="en", slow=False)
            voice.save(f"{assets_dir}/voice_{i}.mp3")
            audio = AudioFileClip(f"{assets_dir}/voice_{i}.mp3")
            img_clip = ImageClip(
                f"{assets_dir}/slide_{i}.png", duration=audio.duration + 1
            )
            video = img_clip.set_audio(audio)
            videos.append(video)
        return videos

    def combine_videos(self, videos):
        """Combines all the videos into one video."""
        final_clip = concatenate_videoclips(videos, method="compose")
        final_clip.write_videofile(self.output_file, fps=24)

    def convert(self):
        """Converts the PowerPoint presentation to a video."""
        self.convert_to_pdf()
        videos = self.create_videos()
        self.combine_videos(videos)


def main():
    """Main function to test the PresentationAutomator class."""
    PresentationAutomator(
        pptx_filename="Final-Project-Team-2.pptx",
        output_file="Final-Project-Team-2.mp4",
    ).convert()


if __name__ == "__main__":
    main()
