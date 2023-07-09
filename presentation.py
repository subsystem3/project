import os
import tempfile
from pptx import Presentation
from gtts import gTTS
from pdf2image import convert_from_path
from moviepy.editor import concatenate_videoclips, VideoFileClip
import ffmpeg
import traceback


def create_pptx(pptx_path: str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Hello, World!"

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Hello, World!"

    prs.save(pptx_path)


def convert_to_pdf(pptx_path: str):
    command = f'unoconv -f pdf -e PageRange=1-2 "{pptx_path}"'
    os.system(command)


def ppt_presenter(pptx_path, pdf_path, output_path):
    with tempfile.TemporaryDirectory() as temp_path:
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        video_list = []
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                tts = gTTS(text=notes, lang="en")
                image_path = os.path.join(temp_path, f"frame_{i}.jpg")
                audio_path = os.path.join(temp_path, f"frame_{i}.mp3")

                image.save(image_path)
                tts.save(audio_path)

                video_path = os.path.join(temp_path, f"frame_{i}.ts")
                ffmpeg_call(image_path, audio_path, temp_path, i)
                if os.path.exists(video_path):
                    video_list.append(video_path)
                else:
                    print(f"Video file {video_path} could not be created")

        ffmpeg_concat(video_list, output_path)


def ffmpeg_call(image_path: str, audio_path: str, temp_path: str, i: int):
    # Check if image_path exists
    if not os.path.exists(image_path):
        print(f"Image file {image_path} does not exist")
        return

    # Check if audio_path exists
    if not os.path.exists(audio_path):
        print(f"Audio file {audio_path} does not exist")
        return

    out_path_mp4 = os.path.join(temp_path, f"frame{i}.mp4")
    out_path_ts = os.path.join(temp_path, f"frame{i}.ts")
    input_still = ffmpeg.input(image_path)
    input_audio = ffmpeg.input(audio_path)
    try:
        ffmpeg.concat(input_still, input_audio, v=1, a=1).output(out_path_mp4).run(
            overwrite_output=True
        )
        # Convert the .mp4 file to .ts
        ffmpeg.input(out_path_mp4).output(out_path_ts, format="mpegts").run()
    except Exception:
        print(traceback.format_exc())
        raise ValueError("A very specific bad thing happened.")


def ffmpeg_concat(video_list, out_path):
    vids = []
    for video_path in video_list:
        if os.path.exists(video_path):
            vids.append(VideoFileClip(video_path))
        else:
            print(f"File {video_path} does not exist")
    final = concatenate_videoclips(vids)
    final.write_videofile(out_path)


def main():
    pptx_path = "test.pptx"
    pdf_path = "test.pdf"
    output_path = "output.mp4"
    create_pptx(pptx_path)
    convert_to_pdf(pptx_path)
    ppt_presenter(pptx_path, pdf_path, output_path)


if __name__ == "__main__":
    main()
