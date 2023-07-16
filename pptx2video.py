import argparse
import hashlib
import os
import shutil
import subprocess
from typing import List

import numpy as np
from google.cloud import texttospeech as tts
from gtts import gTTS
from moviepy.audio.AudioClip import AudioArrayClip
from moviepy.editor import (
    AudioFileClip,
    ImageClip,
    VideoFileClip,
    concatenate_audioclips,
    concatenate_videoclips,
)
from moviepy.video.compositing.transitions import crossfadein, crossfadeout
from pdf2image import convert_from_path
from pptx import Presentation


class PPTXtoVideo:
    """
    A class to automate the creation of a video presentation from a PowerPoint file.
    """

    def __init__(self, pptx_filename: str):
        self.pptx_filename = pptx_filename
        self.pdf_filename = pptx_filename.replace(".pptx", ".pdf")
        self.output_file = pptx_filename.replace(".pptx", ".mp4")
        self.presentation = Presentation(pptx_filename)
        self.slides = self.presentation.slides
        self.voiceover_texts = [
            slide.notes_slide.notes_text_frame.text for slide in self.slides
        ]

    def hash_file(self, filename):
        """
        Calculates the MD5 hash of a file.

        Args:
            filename (str): The name of the file to hash.

        Returns:
            str: The MD5 hash of the file.
        """
        with open(filename, "rb") as f:
            file_hash = hashlib.md5(f.read()).hexdigest()
        return file_hash

    def text_to_wav(self, text: str, filename: str, voice_name: str = "en-US-Studio-M"):
        """
        Converts the given text to speech and saves it as a .wav file.

        If the GOOGLE_APPLICATION_CREDENTIALS environment variable is set, this method uses
        Google Cloud Text-to-Speech to generate the speech. Otherwise, it uses gTTS.

        Args:
            text (str): The text to convert to speech.
            filename (str): The name of the .wav file to save the speech to.
            voice_name (str, optional): The name of the voice to use for speech generation.
                This should be a voice name from Google Cloud Text-to-Speech (e.g., "en-US-Studio-M").
                Defaults to "en-US-Studio-M".
                List of voices at: https://cloud.google.com/text-to-speech/docs/voices
        """
        # USE PROFESSIONAL VOICES FROM GOOGLE CLOUD
        if "GOOGLE_APPLICATION_CREDENTIALS" in os.environ:
            language_code = "-".join(voice_name.split("-")[:2])
            text_input = tts.SynthesisInput(text=text)
            voice_params = tts.VoiceSelectionParams(
                language_code=language_code, name=voice_name
            )
            audio_config = tts.AudioConfig(audio_encoding=tts.AudioEncoding.LINEAR16)
            client = tts.TextToSpeechClient()

            response = client.synthesize_speech(
                input=text_input,
                voice=voice_params,
                audio_config=audio_config,
            )

            with open(filename, "wb") as out:
                out.write(response.audio_content)
        # USE FREE NON-PROFESSIONAL VOICES FROM GTTS
        else:
            voice = gTTS(text=text, lang="en", slow=False)
            voice.save(filename)

    def format_duration(self, duration: int) -> str:
        """
        Formats a duration in seconds into a string in the format 'mm:ss'.

        Args:
            duration (int): Duration in seconds.

        Returns:
            str: Formatted duration string.
        """
        minutes, seconds = divmod(int(duration), 60)
        return f"{minutes}:{seconds:02}"

    def write_metadata(self, videos: List[AudioFileClip]):
        """
        Writes metadata to a text file.

        Args:
            videos (List[AudioFileClip]): List of video clips.
        """
        total_duration = sum(video.duration for video in videos)
        with open(self.pptx_filename.replace(".pptx", ".txt"), "w") as f:
            f.write(f"Total duration: {self.format_duration(total_duration)}\n")
            for i, video in enumerate(videos):
                f.write(f"\nSlide {i+1}:\n")
                f.write(f"Duration: {self.format_duration(video.duration)}\n")
                f.write(f"Voiceover: {self.voiceover_texts[i]}\n")

    def convert_to_pdf(self):
        """
        Converts the .pptx file to a .pdf file using LibreOffice.
        """
        cmd = f"libreoffice --headless --convert-to pdf {self.pptx_filename}"
        subprocess.run(cmd, shell=True, check=True, env={"PATH": "/usr/bin"})


    def create_videos(self, voice_name: str = "en-US-Studio-M") -> List[AudioFileClip]:
        """
        Creates a video for each slide with a voiceover.

        Args:
            voice_name (str, optional): The name of the voice to use for speech generation.
                This should be a voice name from Google Cloud Text-to-Speech (e.g., "en-US-Studio-M").
                Defaults to "en-US-Studio-M".
                List of voices at: https://cloud.google.com/text-to-speech/docs/voices

        Returns:
            List[AudioFileClip]: List of video clips.
        """
        videos = []
        assets_dir = "assets"
        os.makedirs(assets_dir, exist_ok=True)
        for i, slide in enumerate(self.slides):
            text = self.voiceover_texts[i]
            images = convert_from_path(self.pdf_filename, dpi=300)
            image_filename = f"{assets_dir}/slide_{i}.png"
            images[i].save(image_filename, "PNG")

            print(f"Slide {i+1} image saved as {image_filename}")

            # CALCULATE HASHES FOR VOICEOVER TEXT AND SLIDES
            text_hash = hashlib.md5(text.encode()).hexdigest()
            with open(image_filename, "rb") as f:
                image_hash = hashlib.md5(f.read()).hexdigest()

            # CHECK STORED HASHES TO SEE IF SLIDE HAS CHANGED
            os.makedirs("hashes", exist_ok=True)
            video_filename = f"{assets_dir}/video_{i}.mp4"
            if os.path.exists(f"hashes/hashes_{i}.txt"):
                with open(f"hashes/hashes_{i}.txt", "r") as f:
                    stored_values = f.read().splitlines()
                stored_text_hash, stored_image_hash = stored_values[:2]
                stored_voice_name = stored_values[2] if len(stored_values) > 2 else None
                if (
                    text_hash == stored_text_hash
                    and image_hash == stored_image_hash
                    and voice_name == stored_voice_name
                    and os.path.exists(video_filename)
                ):
                    # SKIP IF NO CHANGES DETECTED
                    print(f"No changes detected for slide {i+1}, skipping...")
                    video = VideoFileClip(video_filename)
                else:
                    # STORE HASHES FOR NEXT TIME
                    with open(f"hashes/hashes_{i}.txt", "w") as f:
                        f.write(f"{text_hash}\n{image_hash}\n{voice_name}")

                    # CREATE VOICEOVER
                    voice_filename = f"{assets_dir}/voice_{i}.wav"
                    self.text_to_wav(text, voice_filename, voice_name)
                    print(f"Voiceover for slide {i+1} saved as {voice_filename}")
                    audio = AudioFileClip(voice_filename)

                    # ADD 0.5s SILENCE AT START AND END OF AUDIO (1s TOTAL BETWEEN SLIDES)
                    silence = AudioArrayClip(np.array([[0], [0]]), fps=44100).set_duration(0.5)
                    audio = concatenate_audioclips([silence, audio, silence])

                    # CREATE VIDEO CLIP FROM IMAGE AND AUDIO
                    img_clip = ImageClip(image_filename, duration=audio.duration)
                    img_clip.resize(height=1080)
                    video = img_clip.set_audio(audio)

                    # SAVE EACH VIDEO CLIP
                    video.write_videofile(f"{assets_dir}/video_{i}.mp4", fps=24)
                    print(f"Video for slide {i+1} saved as {video_filename}")

            videos.append(video)

        return videos

    def combine_videos(self, videos: List[AudioFileClip]):
        """
        Combines all the videos into one video.

        Args:
            videos (List[AudioFileClip]): List of video clips.
        """
        intro_clip = VideoFileClip("stock/intro.mp4")
        candidate_intro_hash = self.hash_file("stock/intro.mp4")

        # CHECK HASH FILE FOR INTRO VIDEO
        if os.path.exists("hashes/hash_intro.txt"):
            with open("hashes/hash_intro.txt", "r") as f:
                stored_intro_hash = f.read().strip()

            # CHECK IF ALL INDIVIDUAL SLIDE VIDEOS EXIST
            all_videos_exist = all(os.path.exists(f"assets/video_{i}.mp4") for i in range(len(self.slides)))

            # SKIP VIDEO GENERATION WHEN NO CHANGES DETECTED, FINAL VIDEO EXISTS, AND ALL INDIVIDUAL VIDEOS EXIST
            if (
                all(video is None for video in videos)
                and candidate_intro_hash == stored_intro_hash
                and os.path.exists(self.output_file)
                and all_videos_exist
            ):
                print("No changes detected, skipping video generation...")
                return
        else:
            # STORE HASH FOR INTRO VIDEO
            stored_intro_hash = self.hash_file("stock/intro.mp4")
            with open("hashes/hash_intro.txt", "w") as f:
                f.write(stored_intro_hash)

        intro_clip = crossfadeout(intro_clip, 1)
        videos[0] = crossfadein(videos[0], 1)
        videos.insert(0, intro_clip)
        final_clip = concatenate_videoclips(videos)
        final_clip.write_videofile(self.output_file, fps=24)

    def convert(self):
        self.convert_to_pdf()
        videos = self.create_videos()
        self.write_metadata(videos)
        self.combine_videos(videos)

        # STORE HASH FOR INTRO VIDEO
        self.stored_intro_hash = self.hash_file("stock/intro.mp4")
        with open("hashes/hashes_intro.txt", "w") as f:
            f.write(self.stored_intro_hash)

        # DELETE HASH FILES FOR CONTENT THAT NO LONGER EXISTS
        total_slides = len(self.slides)
        for hash_file in os.scandir("hashes"):
            file_index_str = os.path.splitext(hash_file.name)[0].split("_")[-1]
            if not file_index_str.isdigit():
                continue
            file_index = int(file_index_str)
            if file_index >= total_slides:
                os.remove(hash_file.path)


def main():
    """
    Parse command line args and convert PowerPoint to video.
    """
    parser = argparse.ArgumentParser(
        description="Convert a PowerPoint presentation to a video."
    )
    parser.add_argument(
        "pptx",
        type=str,
        help="The name of the PowerPoint file to convert.",
    )
    parser.add_argument(
        "--keyfile",
        type=str,
        help="The path to the Google service account JSON file.",
        required=False,
    )

    args = parser.parse_args()
    if args.keyfile:
        os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = args.keyfile

    PPTXtoVideo(args.pptx).convert()


if __name__ == "__main__":
    main()
